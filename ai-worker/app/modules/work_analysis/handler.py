"""作品分析核心处理器"""

from __future__ import annotations

import asyncio
import logging
import mimetypes
import tempfile
import threading
from datetime import datetime, timezone
from pathlib import Path
from typing import Callable

from app.config import WorkerSettings, get_settings
from app.providers.openai_compatible import OpenAICompatibleProvider
from app.providers.base import ProviderDescriptor, ProviderType

from .audio_analyzer import AudioFeatures, extract_audio, analyze_audio
from .ai_analyzer import (
    AnalysisResult,
    describe_keyframes_with_ai,
    analyze_comprehensive,
    evaluate_with_criteria,
)
from .criteria import get_default_criteria, parse_criteria_file
from .frame_extractor import Keyframe, extract_keyframes, extract_keyframes_ffmpeg, encode_image_base64
from .metadata import VideoMetadata, extract_metadata
from .ocr import ocr_keyframes

# 音乐分析模块（可选依赖：librosa）
try:
    from .music_analyzer import analyze_music, build_music_analysis_prompt, MusicFeatures
    HAS_MUSIC_ANALYZER = True
except ImportError:
    HAS_MUSIC_ANALYZER = False
    logger.warning("music_analyzer not available. Install librosa for music analysis support.")
from .schemas import (
    AudioAnalysis,
    AudioGuidanceRequest,
    AudioGuidanceResult,
    AudioSegment,
    ContentAnalysis,
    EvaluationResult,
    GuidanceContent,
    KeyframeInfo,
    ScoreItem,
    TechnicalQuality,
    WorkAnalysisOptions,
    WorkAnalysisProgress,
    WorkAnalysisRequest,
    WorkAnalysisResult,
    VideoMetadata as VideoMetadataSchema,
    VideoTaskStatus,
)
from .transcriber import (
    TranscriptionResult,
    transcribe_with_whisper_api,
    transcribe_with_local_whisper,
    transcribe_with_faster_whisper,
    build_speech_analysis,
)

logger = logging.getLogger(__name__)


# 文件类型定义
FILE_TYPE_EXTENSIONS = {
    "video": {".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv", ".m4v", ".3gp"},
    "audio": {".mp3", ".wav", ".flac", ".aac", ".ogg", ".wma", ".m4a", ".opus"},
    "document": {".pdf", ".doc", ".docx", ".txt", ".md", ".ppt", ".pptx", ".xls", ".xlsx", ".rtf", ".odt"},
    "image": {".jpg", ".jpeg", ".png", ".webp", ".bmp"},
}


def detect_file_type(file_path: str) -> str:
    """根据文件后缀检测文件类型

    Args:
        file_path: 文件路径

    Returns:
        文件类型: "video", "audio", "document"

    Raises:
        ValueError: 不支持的文件类型
    """
    ext = Path(file_path).suffix.lower()

    for file_type, extensions in FILE_TYPE_EXTENSIONS.items():
        if ext in extensions:
            return file_type

    raise ValueError(f"不支持的文件类型: {ext}")


def sanitize_file_name(file_name: str) -> str:
    """处理文件名中的特殊字符

    Args:
        file_name: 原始文件名

    Returns:
        处理后的文件名
    """
    import re
    # 替换危险字符
    return re.sub(r'[/\\:*?"<>|]', '_', file_name)


class VideoAnalyzer:
    """视频分析器"""

    def __init__(
        self,
        settings: WorkerSettings | None = None,
        vision_provider: OpenAICompatibleProvider | None = None,
        speech_provider: OpenAICompatibleProvider | None = None,
        text_provider: OpenAICompatibleProvider | None = None,
    ):
        self.settings = settings or get_settings()
        self.vision_provider = vision_provider
        self.speech_provider = speech_provider
        self.text_provider = text_provider
        self.ocr_slots = threading.BoundedSemaphore(self.settings.max_concurrent_ocr_tasks)

        # 临时目录
        self.temp_dir = Path(tempfile.mkdtemp(prefix="work_analysis_"))

        # MinIO 客户端
        self._minio_client = None
        if self.settings.minio_endpoint and self.settings.minio_access_key and self.settings.minio_secret_key:
            try:
                from minio import Minio
                endpoint = self.settings.minio_endpoint.replace("http://", "").replace("https://", "")
                self._minio_client = Minio(
                    endpoint,
                    access_key=self.settings.minio_access_key,
                    secret_key=self.settings.minio_secret_key,
                    secure=self.settings.minio_secure,
                )
            except Exception as e:
                logger.warning("Failed to initialize MinIO client: %s", e)

    async def analyze(
        self,
        request: WorkAnalysisRequest,
        progress_callback: Callable[[WorkAnalysisProgress], None] | None = None,
    ) -> WorkAnalysisResult:
        """执行作品分析"""
        task_id = request.task_id
        started_at = datetime.now(timezone.utc)

        def update_progress(status: VideoTaskStatus, progress: float, detail: str = ""):
            if progress_callback:
                progress_callback(WorkAnalysisProgress(
                    task_id=task_id,
                    status=status,
                    progress=progress,
                    current_stage=status.value,
                    detail=detail,
                ))

        try:
            # 初始化结果
            result = WorkAnalysisResult(
                task_id=task_id,
                file_name=request.file_name,
                status=VideoTaskStatus.PREPROCESSING,
                started_at=started_at.isoformat(),
            )

            # 检测文件类型
            try:
                file_type = detect_file_type(request.file_path)
                result.file_type = file_type
                logger.info("Detected file type: %s for file: %s", file_type, request.file_name)
            except ValueError as e:
                result.status = VideoTaskStatus.FAILED
                result.error = str(e)
                result.completed_at = datetime.now(timezone.utc).isoformat()
                return result

            # 根据文件类型分发到不同的分析流程
            if file_type == "video":
                return await self._analyze_video(request, result, update_progress, started_at)
            elif file_type == "audio":
                return await self._analyze_audio(request, result, update_progress, started_at)
            elif file_type == "document":
                return await self._analyze_document(request, result, update_progress, started_at)
            elif file_type == "image":
                return await self._analyze_images(request, result, update_progress, started_at)
            else:
                result.status = VideoTaskStatus.FAILED
                result.error = f"未实现的文件类型分析: {file_type}"
                result.completed_at = datetime.now(timezone.utc).isoformat()
                return result

        except Exception as e:
            logger.exception("Work analysis failed")
            result.status = VideoTaskStatus.FAILED
            result.error = str(e)
            result.completed_at = datetime.now(timezone.utc).isoformat()
            update_progress(VideoTaskStatus.FAILED, 0, f"处理失败: {e}")
            return result

    async def _analyze_video(
        self,
        request: WorkAnalysisRequest,
        result: WorkAnalysisResult,
        update_progress: Callable,
        started_at: datetime,
    ) -> WorkAnalysisResult:
        """分析视频类型作品"""
        task_id = request.task_id

        # 1. 提取元数据
        update_progress(VideoTaskStatus.EXTRACTING_METADATA, 10, "提取视频元数据")
        metadata = extract_metadata(request.file_path)
        result.metadata = VideoMetadataSchema(
            duration_seconds=metadata.duration,
            width=metadata.width,
            height=metadata.height,
            fps=metadata.fps,
            codec=metadata.work_codec,
            bitrate=0,
            file_size=metadata.file_size,
            format_name=metadata.format_name,
            has_audio=metadata.audio_codec != "unknown",
            audio_codec=metadata.audio_codec,
            audio_sample_rate=metadata.sample_rate,
        )

        # 检查视频时长限制
        if metadata.duration > 30 * 60:
            result.warnings.append("视频超过30分钟，仅分析前30分钟")

        # 2. 提取关键帧
        keyframes_dir = self.temp_dir / task_id / "keyframes"
        keyframes: list[Keyframe] = []

        if request.options.extract_keyframes:
            update_progress(VideoTaskStatus.EXTRACTING_KEYFRAMES, 25, "提取关键帧")

            method = request.options.keyframe_method.value if hasattr(request.options.keyframe_method, 'value') else str(request.options.keyframe_method)
            keyframes = extract_keyframes_ffmpeg(
                request.file_path,
                keyframes_dir,
                method=method,
                threshold=request.options.scene_threshold,
                max_frames=request.options.max_keyframes,
                interval=request.options.min_interval_seconds,
            )

            # 转换关键帧为 base64（确保浏览器可访问）
            for kf in keyframes:
                if kf.path and Path(kf.path).exists():
                    # 直接转换为 base64，不使用 MinIO
                    kf.image_base64 = "data:image/jpeg;base64," + encode_image_base64(kf.path)

            result.keyframes = [
                KeyframeInfo(
                    frame_id=f"frame_{kf.index}",
                    timestamp_seconds=kf.timestamp,
                    frame_index=kf.index,
                    scene_change_score=kf.change_score,
                    image_path=kf.url,       # MinIO URL（前端展示用）
                    image_base64=kf.image_base64,
                )
                for kf in keyframes
            ]
            result.status = VideoTaskStatus.EXTRACTING_KEYFRAMES

        # 3. 提取和分析音频
        audio_path = None
        audio_features = None
        transcription = None

        if request.options.transcribe_audio and metadata.audio_codec != "unknown":
            update_progress(VideoTaskStatus.EXTRACTING_AUDIO, 40, "提取音频")
            audio_dir = self.temp_dir / task_id
            audio_path = extract_audio(request.file_path, audio_dir)

            update_progress(VideoTaskStatus.TRANSCRIBING, 50, "分析音频特征")
            audio_features = analyze_audio(audio_path, metadata.duration)

            # 语音转录
            update_progress(VideoTaskStatus.TRANSCRIBING, 55, "语音转录")
            # 优先使用本地 Whisper（已有缓存模型）
            transcription = transcribe_with_local_whisper(audio_path)
            if not transcription:
                # 尝试 faster-whisper
                transcription = transcribe_with_faster_whisper(audio_path)
            if not transcription:
                # 尝试 Whisper API
                if self.settings.model_api_key and self.settings.model_api_base_url:
                    transcription = transcribe_with_whisper_api(
                        audio_path,
                        self.settings.model_api_key,
                        self.settings.model_api_base_url,
                    )

            # 构建音频分析结果
            if transcription:
                if not transcription.reliable:
                    result.warnings.append(transcription.quality_warning)
                result.audio_analysis = AudioAnalysis(
                    transcription=[
                        AudioSegment(
                            start_time=seg.start_time,
                            end_time=seg.end_time,
                            text=seg.text,
                            confidence=seg.confidence,
                        )
                        for seg in transcription.segments
                    ],
                    total_speech_duration=audio_features.total_speech_duration if audio_features else 0,
                    average_speech_rate=transcription.speech_rate,
                    detected_language=transcription.language,
                )
            else:
                result.audio_analysis = AudioAnalysis(
                    transcription=[],
                    total_speech_duration=audio_features.total_speech_duration if audio_features else 0,
                    average_speech_rate=0,
                    detected_language="unknown",
                )

            result.status = VideoTaskStatus.TRANSCRIBING

        # 4. OCR 关键帧（可选，如果 PaddleOCR 未安装则跳过）
        if request.options.ocr_enabled and keyframes:
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 65, "OCR 文字识别")
            try:
                keyframes = self._run_ocr(keyframes)
            except Exception as e:
                logger.warning("OCR failed, continuing without OCR: %s", e)
                # 设置默认值
                for kf in keyframes:
                    if not kf.ocr_texts:
                        kf.ocr_texts = []
                    if not kf.ocr_summary:
                        kf.ocr_summary = "(OCR not available)"

        # 5. AI 综合分析
        scene_description = ""
        if self.text_provider and keyframes:
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 70, "AI 画面分析")

            # 使用 vision provider 进行画面描述
            client = self.vision_provider.create_client() if self.vision_provider else self.text_provider.create_client()
            model = self.settings.vision_model_name or self.settings.text_model_name or "gpt-5.5"

            scene_description = describe_keyframes_with_ai(client, model, keyframes)

        # 6. AI 综合分析
        ai_analysis_text = ""
        if self.text_provider:
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 80, "AI 综合分析")

            client = self.text_provider.create_client()
            model = self.settings.text_model_name or "gpt-5.5"

            ai_analysis_text = analyze_comprehensive(
                client,
                model,
                metadata,
                audio_features,
                transcription,
                keyframes,
                scene_description,
            )

            logger.info("AI analysis text length: %d chars", len(ai_analysis_text) if ai_analysis_text else 0)

            result.content_analysis = ContentAnalysis(
                overall_topic="视频分析",
                summary=ai_analysis_text if ai_analysis_text else "分析完成",
                key_points=[],
                keywords=[],
            )

        # 7. 评判标准评分（上传文件优先，否则使用默认标准）
        if self.text_provider and ai_analysis_text:
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 90, "评判标准评分")

            client = self.text_provider.create_client()
            model = self.settings.text_model_name or "gpt-5.5"

            criteria_text = request.criteria_text or get_default_criteria()
            logger.info("Using criteria: %s", "uploaded file" if request.criteria_text else "default")

            evaluation_data = evaluate_with_criteria(
                client,
                model,
                ai_analysis_text,
                criteria_text,
                request.supporting_document_text,
                request.supporting_document_name,
            )

            # 将评分结果存入 content_analysis.evaluation
            if result.content_analysis and isinstance(evaluation_data, dict):
                try:
                    result.content_analysis.evaluation = EvaluationResult(
                        total_score=evaluation_data.get("total_score", 0),
                        grade=evaluation_data.get("grade", ""),
                        scores=[
                            ScoreItem(
                                dimension=s.get("dimension", ""),
                                max_score=s.get("max_score", 0),
                                score=s.get("score", 0),
                                evidence=s.get("evidence", ""),
                                suggestion=s.get("suggestion", ""),
                            )
                            for s in evaluation_data.get("scores", [])
                        ],
                        strengths=evaluation_data.get("strengths", []),
                        weaknesses=evaluation_data.get("weaknesses", []),
                        priority_suggestions=evaluation_data.get("priority_suggestions", []),
                        criteria_text=evaluation_data.get("criteria_text", criteria_text),
                        brief_comment=evaluation_data.get("brief_comment", ""),
                        notes=evaluation_data.get("notes", []),
                        document_conformity=evaluation_data.get("document_conformity"),
                    )
                    logger.info("Evaluation stored: total_score=%.1f, grade=%s",
                                result.content_analysis.evaluation.total_score,
                                result.content_analysis.evaluation.grade)
                except Exception as e:
                    logger.error("Failed to store evaluation result: %s", e)

        # 8. 技术质量评估
        result.technical_quality = self._assess_technical_quality(metadata, audio_features)

        # 9. 完成
        completed_at = datetime.now(timezone.utc)
        result.status = VideoTaskStatus.COMPLETED
        result.progress = 100
        result.completed_at = completed_at.isoformat()
        result.processing_time_ms = int((completed_at - started_at).total_seconds() * 1000)

        update_progress(VideoTaskStatus.COMPLETED, 100, "分析完成")

        # 清理音频文件
        if audio_path and Path(audio_path).exists():
            Path(audio_path).unlink(missing_ok=True)

        return result

    async def _analyze_images(
        self,
        request: WorkAnalysisRequest,
        result: WorkAnalysisResult,
        update_progress: Callable,
        started_at: datetime,
    ) -> WorkAnalysisResult:
        """Analyze an image collection as one visual work."""
        image_paths = request.image_paths or [request.file_path]
        valid_paths = [path for path in image_paths if Path(path).is_file()]
        if not valid_paths:
            result.status = VideoTaskStatus.FAILED
            result.error = "未找到可分析的图片文件"
            result.completed_at = datetime.now(timezone.utc).isoformat()
            update_progress(VideoTaskStatus.FAILED, 0, result.error)
            return result

        max_images = min(request.options.max_keyframes, 12)
        if len(valid_paths) > max_images:
            result.warnings.append(f"已选择前 {max_images} 张图片进行分析")
            valid_paths = valid_paths[:max_images]

        update_progress(VideoTaskStatus.EXTRACTING_METADATA, 10, "读取图片信息")
        normalized_dir = self.temp_dir / request.task_id / "image-keyframes"
        normalized_dir.mkdir(parents=True, exist_ok=True)
        normalized_paths: list[str] = []
        from app.modules.document_validation.parser import _save_vision_jpeg
        from PIL import Image

        for index, source_path in enumerate(valid_paths, start=1):
            try:
                output_path = normalized_dir / f"image_{index:02d}.jpg"
                with Image.open(source_path) as image:
                    _save_vision_jpeg(image, output_path)
                normalized_paths.append(str(output_path))
            except Exception as exception:
                logger.warning("Failed to normalize image %s: %s", source_path, exception)
        if not normalized_paths:
            result.status = VideoTaskStatus.FAILED
            result.error = "图片预处理失败"
            update_progress(VideoTaskStatus.FAILED, 0, result.error)
            return result
        valid_paths = normalized_paths

        file_size = sum(Path(path).stat().st_size for path in valid_paths)
        metadata = VideoMetadata(
            duration=float(len(valid_paths)),
            width=0,
            height=0,
            fps=0,
            work_codec="image",
            audio_codec="unknown",
            sample_rate=0,
            channels=0,
            file_size=file_size,
            format_name="image-collection",
        )
        result.metadata = VideoMetadataSchema(
            duration_seconds=metadata.duration,
            width=0,
            height=0,
            fps=0,
            codec=metadata.work_codec,
            bitrate=0,
            file_size=file_size,
            format_name=metadata.format_name,
            has_audio=False,
        )

        update_progress(VideoTaskStatus.EXTRACTING_KEYFRAMES, 25, "整理图片帧")
        keyframes = [
            Keyframe(index=index, timestamp=float(index), path=path)
            for index, path in enumerate(valid_paths)
        ]
        for keyframe in keyframes:
            keyframe.image_base64 = (
                "data:image/jpeg;base64," + encode_image_base64(keyframe.path)
            )
        result.keyframes = [
            KeyframeInfo(
                frame_id=f"image_{keyframe.index}",
                timestamp_seconds=keyframe.timestamp,
                frame_index=keyframe.index,
                image_path=keyframe.path,
                image_base64=keyframe.image_base64,
            )
            for keyframe in keyframes
        ]

        if request.options.ocr_enabled:
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 50, "OCR 文字识别")
            try:
                keyframes = self._run_ocr(keyframes)
            except Exception as exception:
                logger.warning("Image OCR failed, continuing without OCR: %s", exception)
                result.warnings.append("图片文字识别失败，已继续进行画面分析")

        scene_description = ""
        if self.text_provider:
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 60, "AI 图片分批分析")
            client = self.vision_provider.create_client() if self.vision_provider else self.text_provider.create_client()
            model = self.settings.vision_model_name or self.settings.text_model_name or "gpt-5.5"
            scene_description = describe_keyframes_with_ai(
                client,
                model,
                keyframes,
                batch_progress=lambda current, total: update_progress(
                    VideoTaskStatus.ANALYZING_CONTENT,
                    60 + current / total * 15,
                    f"分析图片第 {current}/{total} 批",
                ),
            )

        ai_analysis_text = ""
        if self.text_provider:
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 78, "AI 综合分析")
            client = self.text_provider.create_client()
            model = self.settings.text_model_name or "gpt-5.5"
            ai_analysis_text = analyze_comprehensive(
                client,
                model,
                metadata,
                None,
                None,
                keyframes,
                scene_description,
            )
            result.content_analysis = ContentAnalysis(
                overall_topic="图片作品分析",
                summary=ai_analysis_text or "分析完成",
                key_points=[],
                keywords=[],
            )

        if self.text_provider and ai_analysis_text:
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 90, "评分标准评分")
            criteria_text = request.criteria_text or get_default_criteria()
            evaluation_data = evaluate_with_criteria(
                self.text_provider.create_client(),
                self.settings.text_model_name or "gpt-5.5",
                ai_analysis_text,
                criteria_text,
                request.supporting_document_text,
                request.supporting_document_name,
            )
            if result.content_analysis and isinstance(evaluation_data, dict):
                result.content_analysis.evaluation = EvaluationResult(
                    total_score=evaluation_data.get("total_score", 0),
                    grade=evaluation_data.get("grade", ""),
                    scores=[
                        ScoreItem(
                            dimension=score.get("dimension", ""),
                            max_score=score.get("max_score", 0),
                            score=score.get("score", 0),
                            evidence=score.get("evidence", ""),
                            suggestion=score.get("suggestion", ""),
                        )
                        for score in evaluation_data.get("scores", [])
                    ],
                    strengths=evaluation_data.get("strengths", []),
                    weaknesses=evaluation_data.get("weaknesses", []),
                    priority_suggestions=evaluation_data.get("priority_suggestions", []),
                    criteria_text=evaluation_data.get("criteria_text", criteria_text),
                    brief_comment=evaluation_data.get("brief_comment", ""),
                    notes=evaluation_data.get("notes", []),
                    document_conformity=evaluation_data.get("document_conformity"),
                )

        result.technical_quality = TechnicalQuality(
            work_quality="图片作品",
            audio_quality="不适用",
            stability="不适用",
            overall_score=0,
        )
        completed_at = datetime.now(timezone.utc)
        result.status = VideoTaskStatus.COMPLETED
        result.progress = 100
        result.completed_at = completed_at.isoformat()
        result.processing_time_ms = int((completed_at - started_at).total_seconds() * 1000)
        update_progress(VideoTaskStatus.COMPLETED, 100, "分析完成")
        return result

    def _assess_technical_quality(
        self,
        metadata: VideoMetadata,
        audio_features: AudioFeatures | None,
    ) -> TechnicalQuality:
        """评估技术质量"""
        # 视频质量评估
        if metadata.height >= 1080:
            work_quality = "高清"
        elif metadata.height >= 720:
            work_quality = "标清"
        else:
            work_quality = "低清"

        # 音频质量评估
        if audio_features:
            mean_vol = audio_features.mean_volume
            if -20 <= mean_vol <= -10:
                audio_quality = "清晰"
            elif -30 <= mean_vol < -20 or -10 < mean_vol <= -5:
                audio_quality = "一般"
            else:
                audio_quality = "较差"
        else:
            audio_quality = "未知"

        # 稳定性评估（基于帧率）
        if metadata.fps >= 24:
            stability = "稳定"
        elif metadata.fps >= 15:
            stability = "轻微抖动"
        else:
            stability = "严重抖动"

        # 综合评分
        score = 0
        if metadata.height >= 1080:
            score += 40
        elif metadata.height >= 720:
            score += 30
        else:
            score += 20

        if metadata.fps >= 24:
            score += 30
        elif metadata.fps >= 15:
            score += 20
        else:
            score += 10

        if audio_features and -20 <= audio_features.mean_volume <= -10:
            score += 30
        elif audio_features:
            score += 20
        else:
            score += 10

        return TechnicalQuality(
            work_quality=work_quality,
            audio_quality=audio_quality,
            stability=stability,
            overall_score=score,
        )

    async def _analyze_audio(
        self,
        request: WorkAnalysisRequest,
        result: WorkAnalysisResult,
        update_progress: Callable,
        started_at: datetime,
    ) -> WorkAnalysisResult:
        """分析音频类型作品（自动检测音乐/语音）"""
        task_id = request.task_id

        try:
            # 1. 提取音频元数据
            update_progress(VideoTaskStatus.EXTRACTING_METADATA, 10, "提取音频元数据")
            metadata = extract_metadata(request.file_path)
            result.metadata = VideoMetadataSchema(
                duration_seconds=metadata.duration,
                width=0,
                height=0,
                fps=0,
                codec=metadata.work_codec,
                bitrate=0,
                file_size=metadata.file_size,
                format_name=metadata.format_name,
                has_audio=True,
                audio_codec=metadata.audio_codec,
                audio_sample_rate=metadata.sample_rate,
            )

            # 2. 检测音频类型（音乐/语音）
            update_progress(VideoTaskStatus.EXTRACTING_AUDIO, 20, "检测音频类型")
            is_music = self._detect_music_type(request.file_path, metadata)
            audio_type = "音乐" if is_music else "语音"
            logger.info("Detected audio type: %s", audio_type)

            if is_music:
                # 音乐分析流程
                return await self._analyze_music(
                    request, result, update_progress, started_at, metadata
                )
            else:
                # 语音分析流程
                return await self._analyze_speech(
                    request, result, update_progress, started_at, metadata
                )

        except Exception as e:
            logger.exception("Audio analysis failed")
            result.status = VideoTaskStatus.FAILED
            result.error = str(e)
            result.completed_at = datetime.now(timezone.utc).isoformat()
            update_progress(VideoTaskStatus.FAILED, 0, f"处理失败: {e}")
            return result

    def _detect_music_type(self, file_path: str, metadata: VideoMetadata) -> bool:
        """检测音频是否为音乐类型"""
        try:
            # 1. 基于文件名的启发式判断
            file_name = Path(file_path).stem.lower()
            music_keywords = ['music', 'song', 'track', 'instrumental', '伴奏', '音乐', '歌曲',
                              'bgm', 'ost', 'remix', 'mix', 'beat', '旋律']
            speech_keywords = ['speech', 'talk', 'lecture', 'podcast', 'interview', '演讲', '访谈',
                               '录音', '语音', '朗读', '报告']

            for kw in music_keywords:
                if kw in file_name:
                    return True
            for kw in speech_keywords:
                if kw in file_name:
                    return False

            # 2. 基于音频特征的判断
            audio_features = analyze_audio(file_path, metadata.duration)

            # 音乐通常有更高的语音比例和更规律的节奏
            # 如果语音比例很低，可能是音乐
            if audio_features.speech_ratio < 0.3:
                return True

            # 如果有很多短停顿（可能是歌词），可能是音乐
            if audio_features.short_pauses > 20 and audio_features.speech_ratio < 0.5:
                return True

            # 默认尝试语音转录，如果失败则认为是音乐
            transcription = transcribe_with_local_whisper(file_path, model_name="tiny")
            if not transcription or len(transcription.segments) < 3:
                return True

            # 如果转录结果中文字很少，可能是音乐
            total_text = " ".join(seg.text for seg in transcription.segments)
            if len(total_text) < 50:
                return True

            return False

        except Exception as e:
            logger.warning("Music detection failed, defaulting to speech: %s", e)
            return False

    async def _analyze_music(
        self,
        request: WorkAnalysisRequest,
        result: WorkAnalysisResult,
        update_progress: Callable,
        started_at: datetime,
        metadata: VideoMetadata,
    ) -> WorkAnalysisResult:
        """分析音乐作品（含歌词提取）"""
        if not HAS_MUSIC_ANALYZER:
            result.status = VideoTaskStatus.FAILED
            result.error = "音乐分析功能不可用，请安装依赖：pip install librosa soundfile"
            result.completed_at = datetime.now(timezone.utc).isoformat()
            return result

        try:
            # 1. 音乐特征分析
            update_progress(VideoTaskStatus.EXTRACTING_AUDIO, 25, "分析音乐特征")
            music_features = analyze_music(request.file_path)

            # 2. 尝试提取歌词（语音转录）
            update_progress(VideoTaskStatus.TRANSCRIBING, 40, "提取歌词")
            lyrics = ""
            transcription_segments = []

            # 尝试用 Whisper 提取歌词
            transcription = transcribe_with_local_whisper(request.file_path, model_name="tiny")
            if not transcription:
                transcription = transcribe_with_faster_whisper(request.file_path, model_name="tiny")

            if transcription and transcription.segments and transcription.reliable:
                # 过滤掉可能的噪音转录（太短的片段）
                valid_segments = [seg for seg in transcription.segments if len(seg.text.strip()) > 2]
                if valid_segments:
                    lyrics = "\n".join([
                        f"[{seg.start_time:.1f}s - {seg.end_time:.1f}s] {seg.text}"
                        for seg in valid_segments
                    ])
                    transcription_segments = valid_segments
            elif transcription and not transcription.reliable:
                result.warnings.append(transcription.quality_warning)
            logger.info(
                "Extracted lyrics: %d segments, %d chars",
                len(transcription_segments),
                len(lyrics),
            )

            # 构建音频分析结果
            result.audio_analysis = AudioAnalysis(
                transcription=[
                    AudioSegment(
                        start_time=seg.start_time,
                        end_time=seg.end_time,
                        text=seg.text,
                        confidence=seg.confidence,
                    )
                    for seg in transcription_segments
                ],
                total_speech_duration=0,
                average_speech_rate=0,
                detected_language="music",
                clarity_score=None,
            )

            # 3. AI 综合分析（结合旋律和歌词）
            ai_analysis_text = ""
            if self.text_provider:
                update_progress(VideoTaskStatus.ANALYZING_CONTENT, 60, "AI 音乐分析")

                client = self.text_provider.create_client()
                model = self.settings.text_model_name or "gpt-5.5"

                prompt = build_music_analysis_prompt(music_features, {
                    "format_name": metadata.format_name,
                    "duration": metadata.duration,
                }, lyrics)

                try:
                    response = client.chat.completions.create(
                        model=model,
                        messages=[{"role": "user", "content": prompt}],
                        temperature=0.3,
                        max_tokens=1200,
                    )
                    ai_analysis_text = response.choices[0].message.content or ""
                except Exception as e:
                    logger.error("AI music analysis failed: %s", e)

                # 构建音乐特征摘要
                music_summary = f"""音乐特征分析：

节奏信息：
- BPM：{music_features.bpm:.1f}
- 节拍数：{music_features.beat_count}
- 节奏稳定性：{music_features.tempo_stability:.1%}

调性信息：
- 主调：{music_features.key}
- 调式：{"大调" if music_features.mode == "major" else "小调"}

情绪特征：
- 能量等级：{music_features.energy_level}
- 情绪效价：{music_features.valence:.1%}
- 唤醒度：{music_features.arousal:.1%}

音乐结构：
"""
                for seg in music_features.segments[:6]:
                    music_summary += f"- {seg['start']:.1f}s - {seg['end']:.1f}s: {seg['label']}\n"

                # 添加歌词部分
                if lyrics:
                    music_summary += f"\n歌词内容：\n{lyrics}\n"

                result.content_analysis = ContentAnalysis(
                    overall_topic="音乐分析",
                    summary=music_summary + "\n\n" + (ai_analysis_text if ai_analysis_text else "分析完成"),
                    key_points=[],
                    keywords=[],
                )

            # 4. 评判标准评分
            if self.text_provider and ai_analysis_text:
                update_progress(VideoTaskStatus.ANALYZING_CONTENT, 80, "评判标准评分")

                client = self.text_provider.create_client()
                model = self.settings.text_model_name or "gpt-5.5"

                criteria_text = request.criteria_text or get_default_criteria()

                evaluation_data = evaluate_with_criteria(
                    client,
                    model,
                    ai_analysis_text,
                    criteria_text,
                    request.supporting_document_text,
                    request.supporting_document_name,
                )

                if result.content_analysis and isinstance(evaluation_data, dict):
                    try:
                        result.content_analysis.evaluation = EvaluationResult(
                            total_score=evaluation_data.get("total_score", 0),
                            grade=evaluation_data.get("grade", ""),
                            scores=[
                                ScoreItem(
                                    dimension=s.get("dimension", ""),
                                    max_score=s.get("max_score", 0),
                                    score=s.get("score", 0),
                                    evidence=s.get("evidence", ""),
                                    suggestion=s.get("suggestion", ""),
                                )
                                for s in evaluation_data.get("scores", [])
                            ],
                            strengths=evaluation_data.get("strengths", []),
                            weaknesses=evaluation_data.get("weaknesses", []),
                            priority_suggestions=evaluation_data.get("priority_suggestions", []),
                            criteria_text=evaluation_data.get("criteria_text", criteria_text),
                            brief_comment=evaluation_data.get("brief_comment", ""),
                            notes=evaluation_data.get("notes", []),
                            document_conformity=evaluation_data.get("document_conformity"),
                        )
                    except Exception as e:
                        logger.error("Failed to store evaluation result: %s", e)

            # 5. 完成
            completed_at = datetime.now(timezone.utc)
            result.status = VideoTaskStatus.COMPLETED
            result.progress = 100
            result.completed_at = completed_at.isoformat()
            result.processing_time_ms = int((completed_at - started_at).total_seconds() * 1000)

            update_progress(VideoTaskStatus.COMPLETED, 100, "分析完成")

            return result

        except Exception as e:
            logger.exception("Music analysis failed")
            result.status = VideoTaskStatus.FAILED
            result.error = str(e)
            result.completed_at = datetime.now(timezone.utc).isoformat()
            update_progress(VideoTaskStatus.FAILED, 0, f"处理失败: {e}")
            return result

    async def _analyze_speech(
        self,
        request: WorkAnalysisRequest,
        result: WorkAnalysisResult,
        update_progress: Callable,
        started_at: datetime,
        metadata: VideoMetadata,
    ) -> WorkAnalysisResult:
        """分析语音音频作品"""
        try:
            # 1. 音频特征分析
            update_progress(VideoTaskStatus.EXTRACTING_AUDIO, 25, "分析音频特征")
            audio_features = analyze_audio(request.file_path, metadata.duration)

            # 2. 语音转录
            update_progress(VideoTaskStatus.TRANSCRIBING, 40, "语音转录")
            transcription = transcribe_with_local_whisper(request.file_path)
            if not transcription:
                transcription = transcribe_with_faster_whisper(request.file_path)
            if not transcription:
                if self.settings.model_api_key and self.settings.model_api_base_url:
                    transcription = transcribe_with_whisper_api(
                        request.file_path,
                        self.settings.model_api_key,
                        self.settings.model_api_base_url,
                    )

            # 构建音频分析结果
            if transcription:
                if not transcription.reliable:
                    result.warnings.append(transcription.quality_warning)
                result.audio_analysis = AudioAnalysis(
                    transcription=[
                        AudioSegment(
                            start_time=seg.start_time,
                            end_time=seg.end_time,
                            text=seg.text,
                            confidence=seg.confidence,
                        )
                        for seg in (
                            transcription.segments
                            if transcription.reliable
                            else []
                        )
                    ],
                    total_speech_duration=audio_features.total_speech_duration if audio_features else 0,
                    average_speech_rate=transcription.speech_rate,
                    detected_language=transcription.language,
                    clarity_score=audio_features.clarity_score if audio_features else None,
                )
            else:
                result.audio_analysis = AudioAnalysis(
                    transcription=[],
                    total_speech_duration=audio_features.total_speech_duration if audio_features else 0,
                    average_speech_rate=0,
                    detected_language="unknown",
                )

            # 3. AI 综合分析（基于转录文字）
            ai_analysis_text = ""
            if self.text_provider:
                update_progress(VideoTaskStatus.ANALYZING_CONTENT, 60, "AI 内容分析")

                client = self.text_provider.create_client()
                model = self.settings.text_model_name or "gpt-5.5"

                # 构建音频分析提示
                transcription_text = "\n".join([
                    f"[{seg.start_time:.1f}s - {seg.end_time:.1f}s] {seg.text}"
                    for seg in (
                        transcription.segments
                        if transcription and transcription.reliable
                        else []
                    )
                ])
                if transcription and not transcription.reliable:
                    transcription_text = "ASR 结果可信度过低，已从内容判断和评分证据中排除。"

                prompt = f"""请分析以下音频内容：

音频信息：
- 时长：{metadata.duration:.1f} 秒
- 格式：{metadata.format_name}
- 采样率：{metadata.sample_rate} Hz

转录内容：
{transcription_text}

音频特征：
- 平均语速：{transcription.speech_rate:.1f} 字/分钟（如果可用）
- 检测语言：{transcription.language if transcription else "未知"}

请提供：
1. 主题识别
2. 内容摘要
3. 关键要点
4. 表达质量评估
"""

                try:
                    response = client.chat.completions.create(
                        model=model,
                        messages=[{"role": "user", "content": prompt}],
                        temperature=0.3,
                        max_tokens=1200,
                    )
                    ai_analysis_text = response.choices[0].message.content or ""
                except Exception as e:
                    logger.error("AI analysis failed: %s", e)

                result.content_analysis = ContentAnalysis(
                    overall_topic="音频分析",
                    summary=ai_analysis_text if ai_analysis_text else "分析完成",
                    key_points=[],
                    keywords=[],
                )

            # 4. 评判标准评分
            if self.text_provider and ai_analysis_text:
                update_progress(VideoTaskStatus.ANALYZING_CONTENT, 80, "评判标准评分")

                client = self.text_provider.create_client()
                model = self.settings.text_model_name or "gpt-5.5"

                criteria_text = request.criteria_text or get_default_criteria()

                evaluation_data = evaluate_with_criteria(
                    client,
                    model,
                    ai_analysis_text,
                    criteria_text,
                    request.supporting_document_text,
                    request.supporting_document_name,
                )

                if result.content_analysis and isinstance(evaluation_data, dict):
                    try:
                        result.content_analysis.evaluation = EvaluationResult(
                            total_score=evaluation_data.get("total_score", 0),
                            grade=evaluation_data.get("grade", ""),
                            scores=[
                                ScoreItem(
                                    dimension=s.get("dimension", ""),
                                    max_score=s.get("max_score", 0),
                                    score=s.get("score", 0),
                                    evidence=s.get("evidence", ""),
                                    suggestion=s.get("suggestion", ""),
                                )
                                for s in evaluation_data.get("scores", [])
                            ],
                            strengths=evaluation_data.get("strengths", []),
                            weaknesses=evaluation_data.get("weaknesses", []),
                            priority_suggestions=evaluation_data.get("priority_suggestions", []),
                            criteria_text=evaluation_data.get("criteria_text", criteria_text),
                            brief_comment=evaluation_data.get("brief_comment", ""),
                            notes=evaluation_data.get("notes", []),
                            document_conformity=evaluation_data.get("document_conformity"),
                        )
                    except Exception as e:
                        logger.error("Failed to store evaluation result: %s", e)

            # 5. 完成
            completed_at = datetime.now(timezone.utc)
            result.status = VideoTaskStatus.COMPLETED
            result.progress = 100
            result.completed_at = completed_at.isoformat()
            result.processing_time_ms = int((completed_at - started_at).total_seconds() * 1000)

            update_progress(VideoTaskStatus.COMPLETED, 100, "分析完成")

            return result

        except Exception as e:
            logger.exception("Speech analysis failed")
            result.status = VideoTaskStatus.FAILED
            result.error = str(e)
            result.completed_at = datetime.now(timezone.utc).isoformat()
            update_progress(VideoTaskStatus.FAILED, 0, f"处理失败: {e}")
            return result

    async def _analyze_document(
        self,
        request: WorkAnalysisRequest,
        result: WorkAnalysisResult,
        update_progress: Callable,
        started_at: datetime,
    ) -> WorkAnalysisResult:
        """分析文本文档类型作品"""
        task_id = request.task_id

        try:
            # 1. 提取文档元数据
            update_progress(VideoTaskStatus.EXTRACTING_METADATA, 10, "提取文档元数据")

            file_path = Path(request.file_path)
            file_size = file_path.stat().st_size
            ext = file_path.suffix.lower()
            document_label = "PDF" if ext == ".pdf" else "文档"

            result.metadata = VideoMetadataSchema(
                duration_seconds=0,  # 文档没有时长
                width=0,
                height=0,
                fps=0,
                codec=ext.lstrip('.'),
                bitrate=0,
                file_size=file_size,
                format_name=ext.lstrip('.').upper(),
                has_audio=False,
            )

            # 2. 解析文档内容
            update_progress(VideoTaskStatus.ANALYZING_CONTENT, 25, f"提取{document_label}文字内容")

            document_text = self._parse_document(request.file_path, ext) or ""

            keyframes: list[Keyframe] = []
            image_analysis = ""
            if ext in {".docx", ".pdf"}:
                update_progress(VideoTaskStatus.EXTRACTING_KEYFRAMES, 35, f"提取{document_label}图片关键帧")
                from app.modules.document_validation.parser import extract_docx_images, extract_pdf_images

                try:
                    image_paths = (
                        extract_docx_images(file_path, self.temp_dir / task_id / "document-images")
                        if ext == ".docx"
                        else extract_pdf_images(file_path, self.temp_dir / task_id / "document-images")
                    )
                    keyframes = [
                        Keyframe(index=index, timestamp=float(index), path=str(image_path))
                        for index, image_path in enumerate(image_paths)
                    ]
                    for keyframe in keyframes:
                        keyframe.image_base64 = (
                            f"data:{mimetypes.guess_type(keyframe.path)[0] or 'image/jpeg'};base64,"
                            + encode_image_base64(keyframe.path)
                        )
                    result.keyframes = [
                        KeyframeInfo(
                            frame_id=f"document_image_{keyframe.index}",
                            timestamp_seconds=keyframe.timestamp,
                            frame_index=keyframe.index,
                            image_path=keyframe.path,
                            image_base64=keyframe.image_base64,
                        )
                        for keyframe in keyframes
                    ]

                    if keyframes and request.options.ocr_enabled:
                        update_progress(VideoTaskStatus.ANALYZING_CONTENT, 45, f"OCR 识别{document_label}关键帧")
                        keyframes = self._run_ocr(keyframes)

                    if keyframes and self.text_provider:
                        client = self.vision_provider.create_client() if self.vision_provider else self.text_provider.create_client()
                        model = self.settings.vision_model_name or self.settings.text_model_name or "gpt-5.5"
                        image_analysis = describe_keyframes_with_ai(
                            client,
                            model,
                            keyframes,
                            batch_progress=lambda current, total: update_progress(
                                VideoTaskStatus.ANALYZING_CONTENT,
                                50 + current / total * 20,
                                f"分析{document_label}关键帧第 {current}/{total} 批",
                            ),
                        )
                except Exception as exception:
                    logger.warning("DOCX image extraction failed, continuing with text: %s", exception)
                    result.warnings.append("文档内嵌图片提取失败，已继续进行文字分析")

            if not document_text and not keyframes:
                result.status = VideoTaskStatus.FAILED
                result.error = "无法提取文档文字或图片内容"
                result.completed_at = datetime.now(timezone.utc).isoformat()
                return result

            # 3. AI 综合分析
            ai_analysis_text = ""
            if self.text_provider:
                update_progress(VideoTaskStatus.ANALYZING_CONTENT, 75, f"综合分析{document_label}文字与关键帧")

                client = self.text_provider.create_client()
                model = self.settings.text_model_name or "gpt-5.5"

                # Use a bounded evidence window to reduce latency and token usage.
                text_for_analysis = document_text[:6000] or "该 PDF 未包含可提取正文，请以图片内容为主要分析依据。"

                prompt = f"""请分析以下文档内容：

文档信息：
- 文件名：{request.file_name}
- 文件类型：{ext.lstrip('.').upper()}
- 文件大小：{file_size / 1024:.1f} KB

文档内容：
{text_for_analysis}

文档图片分析：
{image_analysis or "未提取到可分析的内嵌图片"}

请提供：
1. 主题识别
2. 内容摘要
3. 关键要点
4. 文档结构分析
5. 结合正文与图片判断作品内容是否一致
6. 写作质量评估
"""

                try:
                    response = client.chat.completions.create(
                        model=model,
                        messages=[{"role": "user", "content": prompt}],
                        temperature=0.3,
                        max_tokens=1200,
                    )
                    ai_analysis_text = response.choices[0].message.content or ""
                except Exception as e:
                    logger.error("AI analysis failed: %s", e)

                result.content_analysis = ContentAnalysis(
                    overall_topic="文档分析",
                    summary=ai_analysis_text if ai_analysis_text else "分析完成",
                    key_points=[],
                    keywords=[],
                )

            # 4. 评判标准评分
            if self.text_provider and ai_analysis_text:
                update_progress(VideoTaskStatus.ANALYZING_CONTENT, 90, "根据任务标准评分")

                client = self.text_provider.create_client()
                model = self.settings.text_model_name or "gpt-5.5"

                criteria_text = request.criteria_text or get_default_criteria()

                evaluation_data = evaluate_with_criteria(
                    client,
                    model,
                    ai_analysis_text,
                    criteria_text,
                    request.supporting_document_text,
                    request.supporting_document_name,
                )

                if result.content_analysis and isinstance(evaluation_data, dict):
                    try:
                        result.content_analysis.evaluation = EvaluationResult(
                            total_score=evaluation_data.get("total_score", 0),
                            grade=evaluation_data.get("grade", ""),
                            scores=[
                                ScoreItem(
                                    dimension=s.get("dimension", ""),
                                    max_score=s.get("max_score", 0),
                                    score=s.get("score", 0),
                                    evidence=s.get("evidence", ""),
                                    suggestion=s.get("suggestion", ""),
                                )
                                for s in evaluation_data.get("scores", [])
                            ],
                            strengths=evaluation_data.get("strengths", []),
                            weaknesses=evaluation_data.get("weaknesses", []),
                            priority_suggestions=evaluation_data.get("priority_suggestions", []),
                            criteria_text=evaluation_data.get("criteria_text", criteria_text),
                            brief_comment=evaluation_data.get("brief_comment", ""),
                            notes=evaluation_data.get("notes", []),
                            document_conformity=evaluation_data.get("document_conformity"),
                        )
                    except Exception as e:
                        logger.error("Failed to store evaluation result: %s", e)

            # 5. 完成
            completed_at = datetime.now(timezone.utc)
            result.status = VideoTaskStatus.COMPLETED
            result.progress = 100
            result.completed_at = completed_at.isoformat()
            result.processing_time_ms = int((completed_at - started_at).total_seconds() * 1000)

            update_progress(VideoTaskStatus.COMPLETED, 100, "分析完成")

            return result

        except Exception as e:
            logger.exception("Document analysis failed")
            result.status = VideoTaskStatus.FAILED
            result.error = str(e)
            result.completed_at = datetime.now(timezone.utc).isoformat()
            update_progress(VideoTaskStatus.FAILED, 0, f"处理失败: {e}")
            return result

    def _run_ocr(self, keyframes: list[Keyframe]) -> list[Keyframe]:
        """Bound CPU-heavy OCR independently from overall task concurrency."""
        with self.ocr_slots:
            return ocr_keyframes(keyframes, use_paddle=True)

    def _parse_document(self, file_path: str, ext: str) -> str:
        """解析文档内容

        Args:
            file_path: 文件路径
            ext: 文件扩展名

        Returns:
            文档文本内容
        """
        try:
            if ext in {".docx", ".pdf", ".txt", ".md"}:
                from app.modules.document_validation.parser import parse_document

                path = Path(file_path)
                return parse_document(path, path.name, path.stat().st_size).text
            if ext == '.pdf':
                return self._parse_pdf(file_path)
            elif ext in ['.doc', '.docx']:
                return self._parse_docx(file_path)
            elif ext in ['.ppt', '.pptx']:
                return self._parse_pptx(file_path)
            elif ext in ['.xls', '.xlsx']:
                return self._parse_excel(file_path)
            elif ext in ['.txt', '.md', '.rtf']:
                return self._parse_text(file_path)
            else:
                logger.warning("Unsupported document type: %s", ext)
                return ""
        except Exception as e:
            logger.error("Failed to parse document %s: %s", file_path, e)
            return ""

    def _parse_pdf(self, file_path: str) -> str:
        """解析 PDF 文件"""
        try:
            import pdfplumber
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        except ImportError:
            logger.error("pdfplumber not installed. Install with: pip install pdfplumber")
            return ""
        except Exception as e:
            logger.error("Failed to parse PDF: %s", e)
            return ""

    def _parse_docx(self, file_path: str) -> str:
        """解析 Word 文件"""
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs if para.text])
        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return ""
        except Exception as e:
            logger.error("Failed to parse DOCX: %s", e)
            return ""

    def _parse_pptx(self, file_path: str) -> str:
        """解析 PPT 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error("Failed to parse PPTX: %s", e)
            return ""

    def _parse_excel(self, file_path: str) -> str:
        """解析 Excel 文件"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            return text
        except ImportError:
            logger.error("openpyxl not installed. Install with: pip install openpyxl")
            return ""
        except Exception as e:
            logger.error("Failed to parse Excel: %s", e)
            return ""

    def _parse_text(self, file_path: str) -> str:
        """解析纯文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    return f.read()
            except Exception:
                logger.error("Failed to decode text file: %s", file_path)
                return ""
        except Exception as e:
            logger.error("Failed to parse text file: %s", e)
            return ""

    def _upload_to_minio(self, file_path: Path, object_name: str) -> str | None:
        """上传文件到 MinIO 并返回访问 URL"""
        if not self._minio_client:
            return None

        try:
            bucket = self.settings.minio_bucket
            if not self._minio_client.bucket_exists(bucket):
                self._minio_client.make_bucket(bucket)

            self._minio_client.fput_object(
                bucket,
                object_name,
                str(file_path),
                content_type="image/jpeg",
            )

            protocol = "https" if self.settings.minio_secure else "http"
            endpoint = self.settings.minio_endpoint.replace("http://", "").replace("https://", "")
            return f"{protocol}://{endpoint}/{bucket}/{object_name}"
        except Exception as e:
            logger.warning("Failed to upload to MinIO: %s", e)
            return None

    def cleanup(self):
        """清理临时文件"""
        import shutil
        if self.temp_dir.exists():
            shutil.rmtree(self.temp_dir, ignore_errors=True)


# 任务管理器
class WorkAnalysisManager:
    """视频分析任务管理器"""

    def __init__(self):
        self.tasks: dict[str, WorkAnalysisResult] = {}
        self.audio_guidance_tasks: dict[str, AudioGuidanceResult] = {}
        self.analyzer: VideoAnalyzer | None = None
        self.analysis_slots = asyncio.Semaphore(
            get_settings().max_concurrent_work_tasks,
        )

    def get_analyzer(self) -> VideoAnalyzer:
        """获取分析器实例"""
        if not self.analyzer:
            settings = get_settings()

            vision_provider = self._create_provider(settings, ProviderType.VISION, "vision")
            speech_provider = self._create_provider(settings, ProviderType.SPEECH, "speech")
            text_provider = self._create_provider(settings, ProviderType.TEXT, "text")

            self.analyzer = VideoAnalyzer(
                settings=settings,
                vision_provider=vision_provider,
                speech_provider=speech_provider,
                text_provider=text_provider,
            )
        return self.analyzer

    def _create_provider(
        self,
        settings: WorkerSettings,
        provider_type: ProviderType,
        provider_name: str,
    ) -> OpenAICompatibleProvider | None:
        """根据配置创建 provider"""
        if not settings.model_api_base_url or not settings.model_api_key:
            return None

        descriptor = ProviderDescriptor(
            provider_type=provider_type,
            provider_name=f"{provider_name}-provider",
            model_name=self._get_model_name(settings, provider_type),
            base_url=settings.model_api_base_url,
            configured=True,
            required_env_keys=[],
            configured_env_keys=[],
            note=f"{provider_name} model provider",
        )

        return OpenAICompatibleProvider(descriptor, settings)

    def _get_model_name(self, settings: WorkerSettings, provider_type: ProviderType) -> str | None:
        """获取模型名称"""
        model_map = {
            ProviderType.VISION: settings.vision_model_name,
            ProviderType.SPEECH: settings.audio_model_name,
            ProviderType.TEXT: settings.text_model_name,
            ProviderType.MULTIMODAL: settings.multimodal_model_name,
            ProviderType.OCR: settings.ocr_model_name,
            ProviderType.ASR: settings.asr_model_name,
        }
        return model_map.get(provider_type)

    async def submit_task(self, request: WorkAnalysisRequest) -> str:
        """提交分析任务"""
        task_id = request.task_id

        self.tasks[task_id] = WorkAnalysisResult(
            task_id=task_id,
            file_name=request.file_name,
            status=VideoTaskStatus.PENDING,
            progress=0,
        )

        asyncio.create_task(self._execute_task(request))
        return task_id

    async def _execute_task(self, request: WorkAnalysisRequest):
        """执行分析任务"""
        task_id = request.task_id
        analyzer = self.get_analyzer()

        def progress_callback(progress: WorkAnalysisProgress):
            if task_id in self.tasks:
                self.tasks[task_id].status = progress.status
                self.tasks[task_id].progress = progress.progress

        try:
            async with self.analysis_slots:
                result = await asyncio.to_thread(
                    self._run_analysis,
                    analyzer,
                    request,
                    progress_callback,
                )
            self.tasks[task_id] = result
        except Exception as e:
            self.tasks[task_id].status = VideoTaskStatus.FAILED
            self.tasks[task_id].error = str(e)

    @staticmethod
    def _run_analysis(
        analyzer: VideoAnalyzer,
        request: WorkAnalysisRequest,
        progress_callback: Callable[[WorkAnalysisProgress], None],
    ) -> WorkAnalysisResult:
        return asyncio.run(analyzer.analyze(request, progress_callback))

    def get_task(self, task_id: str) -> WorkAnalysisResult | None:
        """获取任务状态"""
        return self.tasks.get(task_id)

    def list_tasks(self) -> list[dict]:
        """列出所有任务"""
        return [
            {
                "task_id": task.task_id,
                "file_name": task.file_name,
                "status": task.status,
                "progress": task.progress,
            }
            for task in self.tasks.values()
        ]


# 全局实例
work_manager = WorkAnalysisManager()
