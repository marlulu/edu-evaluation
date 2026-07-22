from __future__ import annotations

import json
import logging
import re
import subprocess
import tempfile
import wave
import audioop
from contextvars import ContextVar
from uuid import uuid4
from threading import Lock
import zipfile
from pathlib import Path
from collections.abc import Callable
from pydantic import BaseModel, Field
from app.paddle_ocr_compat import create_paddle_ocr, recognize_paddle_text

logger = logging.getLogger(__name__)

_whisper_model = None
_whisper_lock = Lock()
_paddle_ocr = None
_paddle_ocr_lock = Lock()
_progress_callback: ContextVar[Callable[[str], None] | None] = ContextVar(
    "content_analysis_progress_callback",
    default=None,
)
_ARCHIVE_EXTRACT_SUFFIXES = {
    ".avi", ".docx", ".jpeg", ".jpg", ".m4a", ".mkv", ".mov", ".mp3",
    ".mp4", ".pdf", ".png", ".pptx", ".txt", ".wav", ".webp", ".xlsx", ".zip",
    ".css", ".go", ".html", ".java", ".js", ".json", ".jsx", ".py", ".sql",
    ".ts", ".tsx", ".vue", ".xml", ".yaml", ".yml",
}
_CODE_SUFFIXES = {
    ".css", ".go", ".html", ".java", ".js", ".json", ".jsx", ".py", ".sql",
    ".ts", ".tsx", ".vue", ".xml", ".yaml", ".yml",
}


class EvidenceUnit(BaseModel):
    id: str
    file_name: str
    modality: str
    locator: str
    text: str = ""
    metadata: dict[str, object] = Field(default_factory=dict)


class ContentExtractionResult(BaseModel):
    evidence: list[EvidenceUnit] = Field(default_factory=list)
    warnings: list[str] = Field(default_factory=list)
    archive_members: list[dict[str, object]] = Field(default_factory=list)


def extract_content(
        file_paths: list[str],
        progress_callback: Callable[[str], None] | None = None,
) -> ContentExtractionResult:
    result = ContentExtractionResult()
    token = _progress_callback.set(progress_callback)
    try:
        with tempfile.TemporaryDirectory(prefix="submission-extract-") as directory:
            for raw in file_paths:
                _report(f"Starting extraction of {Path(raw).name}")
                _extract(Path(raw), result, None, Path(directory))
    finally:
        _progress_callback.reset(token)
    return result


def _report(message: str) -> None:
    callback = _progress_callback.get()
    if callback is not None:
        callback(message)


def _extract(path: Path, result: ContentExtractionResult, archive_path: str | None, workspace: Path) -> None:
    if not path.is_file():
        result.warnings.append(f"File unavailable: {path}")
        return
    suffix = path.suffix.lower()
    locator_prefix = f"{archive_path}!/" if archive_path else ""
    if suffix in _CODE_SUFFIXES:
        _report(f"Reading source code with line locations from {path.name}")
        _extract_source_code(path, result, locator_prefix)
    elif suffix in {".txt", ".md"}:
        _report(f"Reading structured text from {path.name}")
        _extract_structured_text(path, result, locator_prefix)
    elif suffix == ".pdf":
        _report(f"Parsing PDF pages and embedded images from {path.name}")
        _extract_pdf(path, result, locator_prefix)
    elif suffix == ".docx":
        _report(f"Parsing Word text, tables, and embedded images from {path.name}")
        _extract_docx(path, result, locator_prefix)
    elif suffix == ".pptx":
        _report(f"Parsing PowerPoint slides and embedded images from {path.name}")
        _extract_pptx(path, result, locator_prefix)
    elif suffix == ".xlsx":
        _report(f"Parsing spreadsheet cells, formulas, and images from {path.name}")
        _extract_xlsx(path, result, locator_prefix)
    elif suffix in {".jpg", ".jpeg", ".png", ".webp"}:
        _report(f"Inspecting image metadata for {path.name}")
        evidence_id = f"{path}:image"
        locator = locator_prefix + path.name
        metadata = _image_metadata(path)
        _report(f"Compressing image and preparing Base64 input for {path.name}")
        artifact_key = _upload_artifact(path)
        if artifact_key:
            metadata["artifactObjectKey"] = artifact_key
        result.evidence.append(EvidenceUnit(
            id=evidence_id, file_name=path.name, modality="image", locator=locator, metadata=metadata))
        _extract_image_ocr(path, result, evidence_id, locator)
    elif suffix in {".mp4", ".mov", ".avi", ".mkv", ".mp3", ".wav", ".m4a"}:
        _report(f"Inspecting media streams and duration for {path.name}")
        result.evidence.append(EvidenceUnit(id=f"{path}:media", file_name=path.name, modality="video" if suffix in {".mp4", ".mov", ".avi", ".mkv"} else "audio", locator=locator_prefix + path.name, metadata=_media_metadata(path)))
        if suffix in {".mp4", ".mov", ".avi", ".mkv"}:
            _report(f"Extracting representative keyframes from {path.name}")
            _extract_keyframes(path, result, locator_prefix)
        if suffix in {".mp3", ".wav", ".m4a", ".mp4", ".mov", ".avi", ".mkv"}:
            _report(f"Transcribing audio for {path.name}")
            _transcribe_audio(path, result, locator_prefix)
            _report(f"Measuring audio quality for {path.name}")
            _analyze_audio_quality(path, result, locator_prefix)
        if suffix in {".mp4", ".mov", ".avi", ".mkv"}:
            _build_video_windows(path, result, locator_prefix)
    elif suffix == ".zip":
        try:
            _report(f"Opening archive {path.name} and listing its contents")
            with zipfile.ZipFile(path) as archive:
                members = [member for member in archive.infolist() if not member.is_dir()]
                # Let text and images become visible before a long video transcode/transcription.
                members.sort(key=lambda member: _archive_member_priority(member.filename))
                for index, member in enumerate(members):
                    if index >= 200:
                        result.warnings.append(f"Archive member limit reached: {path.name}")
                        break
                    member_name = _archive_member_name(member)
                    member_path = Path(member_name)
                    if member_path.is_absolute() or ".." in member_path.parts:
                        result.warnings.append(f"Unsafe archive member skipped: {member_name}")
                        continue
                    # ZIP member names are not valid Windows file names by contract.
                    # Never use them for local paths; preserve them only in evidence locators.
                    suffix = member_path.suffix.lower()
                    archive_member = {
                        "archive": path.name,
                        "memberPath": member_name,
                        "extension": suffix,
                        "status": "extracting",
                    }
                    result.archive_members.append(archive_member)
                    _report(f"Extracting archive member {index + 1}/{len(members)}: {member_name}")
                    local_suffix = suffix if suffix in _ARCHIVE_EXTRACT_SUFFIXES else ".bin"
                    extracted = workspace / f"archive-member-{index:03d}-{uuid4().hex}{local_suffix}"
                    try:
                        extracted.write_bytes(archive.read(member))
                        _extract(extracted, result, f"{path.name}!/{member_name}", workspace)
                        archive_member["status"] = "completed"
                        _report(f"Completed archive member {member_name}")
                    except Exception as error:
                        archive_member["status"] = "failed"
                        archive_member["error"] = str(error)[:240]
                        result.warnings.append(
                            f"Archive member extraction failed for {member_name}: {str(error)[:240]}")
        except zipfile.BadZipFile:
            result.warnings.append(f"Unreadable archive: {path.name}")
    else:
        result.warnings.append(f"Unsupported content extraction type: {path.name}")


def _archive_member_name(member: zipfile.ZipInfo) -> str:
    """Recover common GBK ZIP names while retaining UTF-8 names verbatim."""
    if member.flag_bits & 0x800:
        return member.filename
    try:
        return member.filename.encode("cp437").decode("gbk")
    except UnicodeError:
        return member.filename


def _archive_member_priority(member_name: str) -> int:
    return 1 if Path(member_name).suffix.lower() in {
        ".avi", ".m4a", ".mkv", ".mov", ".mp3", ".mp4", ".wav",
    } else 0


def _extract_structured_text(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    text = path.read_text(encoding="utf-8", errors="replace")[:20000]
    result.evidence.append(EvidenceUnit(
        id=f"{path}:text",
        file_name=path.name,
        modality="text-document",
        locator=f"{prefix}{path.name}",
        text=text,
    ))
    sections = _structured_text_sections(text)
    for index, section in enumerate(sections, 1):
        result.evidence.append(EvidenceUnit(
            id=f"{path}:text-section:{index}",
            file_name=path.name,
            modality="text-section",
            locator=f"{prefix}{path.name}#paragraph={section['startParagraph']}-{section['endParagraph']}",
            text=section["text"],
            metadata={
                "heading": section["heading"],
                "headingLevel": section["headingLevel"],
                "startParagraph": section["startParagraph"],
                "endParagraph": section["endParagraph"],
            },
        ))


def _structured_text_sections(text: str) -> list[dict[str, object]]:
    """Split plain and Markdown text into compact, location-aware sections."""
    paragraphs = [value.strip() for value in re.split(r"\n\s*\n", text) if value.strip()]
    sections: list[dict[str, object]] = []
    heading = ""
    heading_level = 0
    section_start = 1
    section_parts: list[str] = []
    heading_pattern = re.compile(r"^(#{1,6})\s+(.+)$|^((?:\d+(?:\.\d+)*[.、)]|[一二三四五六七八九十]+[、.]))\s*(.+)$")

    def append_section(end_paragraph: int) -> None:
        if not section_parts:
            return
        sections.append({
            "heading": heading,
            "headingLevel": heading_level,
            "startParagraph": section_start,
            "endParagraph": end_paragraph,
            "text": "\n\n".join(section_parts)[:6000],
        })

    for paragraph_index, paragraph in enumerate(paragraphs, 1):
        match = heading_pattern.match(paragraph)
        if match:
            append_section(paragraph_index - 1)
            if match.group(1):
                heading = match.group(2).strip()
                heading_level = len(match.group(1))
            else:
                heading = match.group(4).strip()
                heading_level = max(1, match.group(3).count(".") + 1)
            section_start = paragraph_index
            section_parts = [paragraph]
        else:
            section_parts.append(paragraph)
    append_section(len(paragraphs))
    return sections[:80]


def _extract_source_code(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    text = path.read_text(encoding="utf-8", errors="replace")[:80_000]
    lines = text.splitlines()
    if not lines:
        lines = [""]
    for start in range(0, len(lines), 80):
        end = min(start + 80, len(lines))
        result.evidence.append(EvidenceUnit(
            id=f"{path}:source-code:{start + 1}-{end}",
            file_name=path.name,
            modality="source-code",
            locator=f"{prefix}{path.name}#L{start + 1}-L{end}",
            text="\n".join(lines[start:end]),
            metadata={
                "language": path.suffix.lower().lstrip("."),
                "startLine": start + 1,
                "endLine": end,
                "totalLines": len(lines),
            },
        ))
    _report(f"Extracted {len(lines)} source lines from {path.name} into {max(1, (len(lines) + 79) // 80)} code sections")


def _extract_pdf(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for index, page in enumerate(pdf.pages, 1):
                image_ids: list[str] = []
                for image_index, image in enumerate(page.images, 1):
                    image_ids.append(f"{path}:page:{index}:image:{image_index}")
                if image_ids:
                    render_path = path.parent / f"{path.stem}-page-{index}.png"
                    page.to_image(resolution=120).save(str(render_path), format="PNG")
                    page_evidence_id = f"{path}:page:{index}:render"
                    image_ids.append(page_evidence_id)
                    result.evidence.append(EvidenceUnit(id=page_evidence_id, file_name=path.name, modality="pdf-page-image", locator=f"{prefix}{path.name}#page={index}:render", metadata={"artifactObjectKey": _upload_artifact(render_path)}))
                    _extract_image_ocr(render_path, result, page_evidence_id, f"{prefix}{path.name}#page={index}:render")
                tables = page.extract_tables() or []
                for table_index, table in enumerate(tables, 1):
                    table_text = "\n".join("\t".join(cell or "" for cell in row) for row in table)
                    result.evidence.append(EvidenceUnit(
                        id=f"{path}:page:{index}:table:{table_index}",
                        file_name=path.name,
                        modality="pdf-table",
                        locator=f"{prefix}{path.name}#page={index}:table={table_index}",
                        text=table_text[:12000],
                    ))
                result.evidence.append(EvidenceUnit(id=f"{path}:page:{index}", file_name=path.name, modality="pdf-page", locator=f"{prefix}{path.name}#page={index}", text=(page.extract_text() or "")[:12000], metadata={"tableCount": len(tables), "imageEvidenceIds": image_ids, "imageCount": len(image_ids)}))
    except Exception as error:
        result.warnings.append(f"PDF extraction failed for {path.name}: {error}")


def _extract_docx(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    try:
        from docx import Document
        document = Document(path)
        text_parts: list[str] = []
        embedded_ids: list[str] = []
        for index, paragraph in enumerate(document.paragraphs, 1):
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
                result.evidence.append(EvidenceUnit(id=f"{path}:paragraph:{index}", file_name=path.name, modality="word-paragraph", locator=f"{prefix}{path.name}#paragraph={index}", text=paragraph.text))
        for index, table in enumerate(document.tables, 1):
            text = "\n".join("\t".join(cell.text for cell in row.cells) for row in table.rows)
            result.evidence.append(EvidenceUnit(id=f"{path}:table:{index}", file_name=path.name, modality="word-table", locator=f"{prefix}{path.name}#table={index}", text=text[:12000]))
        with zipfile.ZipFile(path) as archive:
            media = [member for member in archive.namelist() if member.startswith("word/media/")]
            for index, member in enumerate(media, 1):
                target = path.parent / f"{path.stem}-embedded-{index}{Path(member).suffix}"
                target.write_bytes(archive.read(member))
                metadata = {"embeddedPath": member, "artifactObjectKey": _upload_artifact(target)}
                evidence_id = f"{path}:embedded:{index}"
                embedded_ids.append(evidence_id)
                result.evidence.append(EvidenceUnit(id=evidence_id, file_name=path.name, modality="word-embedded-image", locator=f"{prefix}{path.name}#{member}", metadata=metadata))
                if target.suffix.lower() in {".jpg", ".jpeg", ".png", ".webp"}:
                    _extract_image_ocr(target, result, evidence_id, f"{prefix}{path.name}#{member}")
        result.evidence.append(EvidenceUnit(id=f"{path}:document-context", file_name=path.name, modality="word-document-context", locator=f"{prefix}{path.name}", text="\n".join(text_parts)[:20000], metadata={"embeddedEvidenceIds": embedded_ids}))
    except Exception as error:
        result.warnings.append(f"Word extraction failed for {path.name}: {error}")


def _extract_pptx(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    try:
        from pptx import Presentation
        presentation = Presentation(path)
        for index, slide in enumerate(presentation.slides, 1):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
            notes = ""
            if slide.has_notes_slide:
                notes = "\n".join(shape.text for shape in slide.notes_slide.shapes if hasattr(shape, "text") and shape.text.strip())
            image_ids: list[str] = []
            for image_index, shape in enumerate((shape for shape in slide.shapes if shape.shape_type == 13), 1):
                target = path.parent / f"{path.stem}-slide-{index}-image-{image_index}.png"
                target.write_bytes(shape.image.blob)
                evidence_id = f"{path}:slide:{index}:image:{image_index}"
                image_ids.append(evidence_id)
                result.evidence.append(EvidenceUnit(id=evidence_id, file_name=path.name, modality="slide-image", locator=f"{prefix}{path.name}#slide={index}:image={image_index}", metadata={"artifactObjectKey": _upload_artifact(target)}))
                _extract_image_ocr(target, result, evidence_id, f"{prefix}{path.name}#slide={index}:image={image_index}")
            result.evidence.append(EvidenceUnit(id=f"{path}:slide:{index}", file_name=path.name, modality="slide", locator=f"{prefix}{path.name}#slide={index}", text=(text + "\nNotes:\n" + notes)[:12000], metadata={"imageCount": len(image_ids), "imageEvidenceIds": image_ids}))
        with zipfile.ZipFile(path) as archive:
            media_members = [member for member in archive.namelist() if member.startswith("ppt/media/")]
            for index, member in enumerate(media_members, 1):
                target = path.parent / f"{path.stem}-media-{index}{Path(member).suffix}"
                target.write_bytes(archive.read(member))
                suffix = target.suffix.lower()
                modality = "slide-embedded-media"
                metadata: dict[str, object] = {"embeddedPath": member}
                if suffix in {".jpg", ".jpeg", ".png", ".webp"}:
                    modality = "slide-image"
                    metadata.update(_image_metadata(target))
                    metadata["artifactObjectKey"] = _upload_artifact(target)
                    _extract_image_ocr(target, result, f"{path}:embedded-media:{index}", f"{prefix}{path.name}#{member}")
                elif suffix in {".mp4", ".mov", ".avi", ".mkv", ".mp3", ".wav", ".m4a"}:
                    modality = "slide-embedded-media"
                    metadata.update(_media_metadata(target))
                result.evidence.append(EvidenceUnit(
                    id=f"{path}:embedded-media:{index}",
                    file_name=path.name,
                    modality=modality,
                    locator=f"{prefix}{path.name}#{member}",
                    metadata=metadata,
                ))
    except Exception as error:
        result.warnings.append(f"PowerPoint extraction failed for {path.name}: {error}")


def _extract_xlsx(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    try:
        import openpyxl
        workbook = openpyxl.load_workbook(path, read_only=False, data_only=False)
        for sheet in workbook.worksheets:
            rows = ["\t".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True)]
            formulas = sum(1 for row in sheet.iter_rows() for cell in row if isinstance(cell.value, str) and cell.value.startswith("="))
            image_ids: list[str] = []
            for image_index, image in enumerate(sheet._images, 1):
                target = path.parent / f"{path.stem}-{sheet.title}-image-{image_index}.{image.format or 'png'}"
                target.write_bytes(image._data())
                evidence_id = f"{path}:sheet:{sheet.title}:image:{image_index}"
                image_ids.append(evidence_id)
                metadata = _image_metadata(target)
                metadata["artifactObjectKey"] = _upload_artifact(target)
                result.evidence.append(EvidenceUnit(
                    id=evidence_id,
                    file_name=path.name,
                    modality="spreadsheet-image",
                    locator=f"{prefix}{path.name}#sheet={sheet.title}:image={image_index}",
                    metadata=metadata,
                ))
                _extract_image_ocr(target, result, evidence_id, f"{prefix}{path.name}#sheet={sheet.title}:image={image_index}")
            result.evidence.append(EvidenceUnit(id=f"{path}:sheet:{sheet.title}", file_name=path.name, modality="spreadsheet", locator=f"{prefix}{path.name}#sheet={sheet.title}", text="\n".join(rows)[:20000], metadata={"formulaCount": formulas, "chartCount": len(sheet._charts), "tableCount": len(sheet.tables), "imageCount": len(image_ids), "imageEvidenceIds": image_ids}))
    except Exception as error:
        result.warnings.append(f"Excel extraction failed for {path.name}: {error}")


def _image_metadata(path: Path) -> dict[str, object]:
    metadata: dict[str, object] = {"sizeBytes": path.stat().st_size}
    try:
        from PIL import Image
        with Image.open(path) as image:
            metadata.update({
                "width": image.width,
                "height": image.height,
                "format": image.format or "",
                "aspectRatio": round(image.width / image.height, 3) if image.height else 0,
            })
            try:
                import cv2
                import numpy as np

                pixels = np.asarray(image.convert("L"))
                brightness = float(pixels.mean())
                contrast = float(pixels.std())
                sharpness = float(cv2.Laplacian(pixels, cv2.CV_64F).var())
                moments = cv2.moments(pixels)
                if moments["m00"]:
                    centroid_x = moments["m10"] / moments["m00"] / image.width
                    centroid_y = moments["m01"] / moments["m00"] / image.height
                else:
                    centroid_x, centroid_y = 0.5, 0.5
                balance = max(0.0, 1 - ((centroid_x - 0.5) ** 2 + (centroid_y - 0.5) ** 2) ** 0.5 * 2)
                metadata.update({
                    "brightness": round(brightness, 1),
                    "contrast": round(contrast, 1),
                    "sharpness": round(sharpness, 1),
                    "compositionBalance": round(balance, 2),
                    "visualSummary": (
                        f"{image.width}x{image.height}, brightness {brightness:.0f}, "
                        f"contrast {contrast:.0f}, sharpness {sharpness:.0f}, "
                        f"balance {balance:.2f}"
                    ),
                })
            except Exception:
                metadata["visualInspectionWarning"] = "Visual quality metrics unavailable"
    except Exception:
        metadata["inspectionWarning"] = "Image dimensions unavailable"
    return metadata


def _extract_image_ocr(
        path: Path,
        result: ContentExtractionResult,
        source_evidence_id: str,
        source_locator: str,
) -> None:
    try:
        _report(f"Running local OCR on {path.name}")
        ocr = _get_paddle_ocr()
        with _paddle_ocr_lock:
            lines = recognize_paddle_text(ocr, str(path))
        text = "\n".join(value for value, confidence in lines if confidence > 0.5)
        if text.strip():
            metadata = {"sourceEvidenceId": source_evidence_id}
            result.evidence.append(EvidenceUnit(
                id=f"{source_evidence_id}:ocr",
                file_name=path.name,
                modality="image-ocr",
                locator=f"{source_locator}#ocr",
                text=text[:12000],
                metadata=metadata,
            ))
            if is_code_like_text(text):
                result.evidence.append(EvidenceUnit(
                    id=f"{source_evidence_id}:code",
                    file_name=path.name,
                    modality="image-code",
                    locator=f"{source_locator}#code",
                    text=text[:12000],
                    metadata=metadata,
                ))
            _report(
                f"OCR completed for {path.name}: {len(lines)} recognized lines; "
                f"text preview: {text[:360].replace(chr(10), ' ')}"
            )
        else:
            _report(f"OCR completed for {path.name}: no readable text")
    except Exception as error:
        result.warnings.append(f"Image OCR unavailable for {path.name}: {error}")


def is_code_like_text(text: str) -> bool:
    normalized = text.strip()
    if len(normalized) < 12:
        return False
    keyword_count = len(re.findall(
        r"\b(?:async|await|class|const|def|else|for|function|if|import|let|return|while)\b",
        normalized,
        flags=re.IGNORECASE,
    ))
    symbol_count = len(re.findall(r"(?:=>|==|!=|[{};]|\w+\([^)]*\)|\b\w+\.\w+\b)", normalized))
    return (keyword_count >= 1 and symbol_count >= 1) or symbol_count >= 3


def _get_paddle_ocr():
    global _paddle_ocr
    if _paddle_ocr is None:
        with _paddle_ocr_lock:
            if _paddle_ocr is None:
                _paddle_ocr = create_paddle_ocr()
    return _paddle_ocr


def _media_metadata(path: Path) -> dict[str, object]:
    metadata: dict[str, object] = {"sizeBytes": path.stat().st_size}
    try:
        completed = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration,format_name", "-show_streams", "-of", "json", str(path)],
            capture_output=True, text=True, timeout=15, check=True,
        )
        payload = json.loads(completed.stdout)
        format_info = payload.get("format", {})
        streams = payload.get("streams", [])
        metadata.update({
            "durationSeconds": float(format_info.get("duration") or 0),
            "format": format_info.get("format_name") or "",
            "videoStreams": sum(stream.get("codec_type") == "video" for stream in streams),
            "audioStreams": sum(stream.get("codec_type") == "audio" for stream in streams),
        })
    except Exception:
        metadata["inspectionWarning"] = "Media metadata unavailable (ffprobe required)"
    return metadata


def _sanitize_filename(name: str) -> str:
    """清理文件名中的特殊字符，仅保留 ASCII 安全字符"""
    import re
    # 替换所有非 ASCII 字符和特殊字符为下划线
    return re.sub(r'[^a-zA-Z0-9._\-]', '_', name)


def _extract_keyframes(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    try:
        import cv2
        capture = cv2.VideoCapture(str(path))
        fps = capture.get(cv2.CAP_PROP_FPS) or 1
        frame_count = capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0
        duration = frame_count / fps
        window_starts = list(range(0, max(1, int(duration)), 60))
        candidates = sorted({
            0,
            *(min(start + 2, max(0, int(duration) - 1)) for start in window_starts),
            *(min(start + 30, max(0, int(duration) - 1)) for start in window_starts),
            max(0, int(duration) - 1),
        })
        # Keep at most two representative frames per 60-second window. Long videos
        # preserve timeline coverage while avoiding an unbounded visual prompt.
        if len(candidates) > 48:
            candidate_indexes = {
                round(index * (len(candidates) - 1) / 47) for index in range(48)
            }
            candidates = [candidate for index, candidate in enumerate(candidates) if index in candidate_indexes]
        selected: list[tuple[int, object]] = []
        for second in candidates:
            capture.set(cv2.CAP_PROP_POS_MSEC, second * 1000)
            ok, frame = capture.read()
            if ok:
                selected.append((second, frame))
        _report(f"Selected {len(selected)} keyframes from {path.name}")
        # 使用安全的文件名（避免特殊字符导致文件保存失败）
        safe_stem = _sanitize_filename(path.stem)
        for index, (second, frame) in enumerate(selected):
            frame_path = path.parent / f"{safe_stem}-frame-{index}.jpg"
            cv2.imwrite(str(frame_path), frame)
            metadata = _image_metadata(frame_path)
            metadata["timestampSeconds"] = second
            _report(f"Compressing keyframe {index + 1}/{len(selected)} at {second}s from {path.name}")
            artifact_key = _upload_artifact(frame_path)
            if artifact_key:
                metadata["artifactObjectKey"] = artifact_key
            evidence_id = f"{path}:frame:{index}"
            locator = f"{prefix}{path.name}#t={second}s"
            result.evidence.append(EvidenceUnit(id=evidence_id, file_name=path.name, modality="video-frame", locator=locator, metadata=metadata))
            _extract_image_ocr(frame_path, result, evidence_id, locator)
        _report(f"Completed keyframe extraction for {path.name}")
        capture.release()
    except Exception as error:
        result.warnings.append(f"Keyframe extraction failed for {path.name}: {error}")


def _transcribe_audio(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    try:
        model = _get_whisper_model()
        if model is None:
            result.warnings.append(f"Audio transcription skipped for {path.name}: no whisper model available")
            return

        # 检测是否为 faster_whisper 模型
        is_faster_whisper = 'WhisperModel' in type(model).__name__

        if is_faster_whisper:
            # faster_whisper API
            segments_iter, _ = model.transcribe(str(path), vad_filter=True)
            values = list(segments_iter)
            for index, segment in enumerate(values):
                result.evidence.append(EvidenceUnit(id=f"{path}:audio:{index}", file_name=path.name, modality="audio-transcript", locator=f"{prefix}{path.name}#t={segment.start:.2f}-{segment.end:.2f}", text=segment.text.strip(), metadata={"startSeconds": segment.start, "endSeconds": segment.end}))
            if values:
                spoken_seconds = sum(max(0, item.end - item.start) for item in values)
                characters = sum(len(item.text.strip()) for item in values)
                pauses = [max(0, values[index + 1].start - item.end) for index, item in enumerate(values[:-1])]
                result.evidence.append(EvidenceUnit(id=f"{path}:speech-rhythm", file_name=path.name, modality="speech-rhythm", locator=f"{prefix}{path.name}#speech", metadata={"charactersPerMinute": round(characters / spoken_seconds * 60, 1) if spoken_seconds else 0, "longPauseCount": sum(pause >= 2 for pause in pauses), "longestPauseSeconds": round(max(pauses, default=0), 2)}))
        else:
            # 标准 whisper API
            result_data = model.transcribe(str(path), language="zh")
            segments = result_data.get("segments", [])
            for index, segment in enumerate(segments):
                start = segment.get("start", 0)
                end = segment.get("end", 0)
                text = segment.get("text", "").strip()
                result.evidence.append(EvidenceUnit(id=f"{path}:audio:{index}", file_name=path.name, modality="audio-transcript", locator=f"{prefix}{path.name}#t={start:.2f}-{end:.2f}", text=text, metadata={"startSeconds": start, "endSeconds": end}))
            if segments:
                spoken_seconds = sum(max(0, seg.get("end", 0) - seg.get("start", 0)) for seg in segments)
                characters = sum(len(seg.get("text", "").strip()) for seg in segments)
                pauses = [max(0, segments[index + 1].get("start", 0) - seg.get("end", 0)) for index, seg in enumerate(segments[:-1])]
                result.evidence.append(EvidenceUnit(id=f"{path}:speech-rhythm", file_name=path.name, modality="speech-rhythm", locator=f"{prefix}{path.name}#speech", metadata={"charactersPerMinute": round(characters / spoken_seconds * 60, 1) if spoken_seconds else 0, "longPauseCount": sum(pause >= 2 for pause in pauses), "longestPauseSeconds": round(max(pauses, default=0), 2)}))
    except Exception as error:
        result.warnings.append(f"Audio transcription failed for {path.name}: {error}")


def _get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        with _whisper_lock:
            if _whisper_model is None:
                try:
                    from faster_whisper import WhisperModel
                    _whisper_model = WhisperModel("base", device="cpu", compute_type="int8")
                    logger.info("Loaded faster_whisper model")
                except ImportError:
                    # 尝试使用标准 whisper
                    try:
                        import whisper
                        _whisper_model = whisper.load_model("base")
                        logger.info("Loaded standard whisper model")
                    except ImportError:
                        logger.warning("Neither faster_whisper nor whisper is installed. Audio transcription will be skipped.")
                        return None
    return _whisper_model


def _analyze_audio_quality(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    try:
        source = path
        temp_wav = None
        if path.suffix.lower() != ".wav":
            # 使用安全的临时文件路径，避免特殊字符问题
            import tempfile
            temp_wav = tempfile.NamedTemporaryFile(suffix=".analysis.wav", delete=False)
            temp_wav.close()
            source = Path(temp_wav.name)
            subprocess.run(["ffmpeg", "-y", "-i", str(path), "-ac", "1", "-ar", "16000", str(source)], capture_output=True, timeout=120, check=True)
        with wave.open(str(source), "rb") as audio:
            rate, width, frames = audio.getframerate(), audio.getsampwidth(), audio.getnframes()
            raw = audio.readframes(frames)
            rms = audioop.rms(raw, width)
            peak = audioop.max(raw, width)
            duration = frames / rate if rate else 0
            result.evidence.append(EvidenceUnit(id=f"{path}:audio-quality", file_name=path.name, modality="audio-quality", locator=f"{prefix}{path.name}#audio", metadata={"durationSeconds": duration, "sampleRate": rate, "rms": rms, "peak": peak}))
        # 清理临时文件
        if temp_wav and source.exists():
            source.unlink(missing_ok=True)
    except Exception as error:
        result.warnings.append(f"Audio quality analysis failed for {path.name}: {error}")


def _build_video_windows(path: Path, result: ContentExtractionResult, prefix: str) -> None:
    media = next((item for item in result.evidence if item.id == f"{path}:media"), None)
    duration = float((media.metadata.get("durationSeconds", 0) if media else 0) or 0)
    for start in range(0, max(1, int(duration)), 60):
        end = min(start + 60, duration)
        evidence_ids: list[str] = []
        frame_ids: list[str] = []
        audio_ids: list[str] = []
        visual_summaries: list[str] = []
        audio_metrics: list[dict[str, object]] = []
        transcript: list[str] = []
        for item in result.evidence:
            if not item.id.startswith(f"{path}:"):
                continue
            timestamp = item.metadata.get("timestampSeconds", item.metadata.get("startSeconds"))
            if isinstance(timestamp, (int, float)) and start <= timestamp < end:
                evidence_ids.append(item.id)
                if item.modality == "video-frame":
                    frame_ids.append(item.id)
                    summary = item.metadata.get("visualSummary")
                    if isinstance(summary, str):
                        visual_summaries.append(summary)
                if item.modality in {"audio-transcript", "audio-quality", "speech-rhythm"}:
                    audio_ids.append(item.id)
                    if item.modality != "audio-transcript":
                        audio_metrics.append(item.metadata)
                if item.modality == "audio-transcript":
                    transcript.append(item.text)
            elif item.modality in {"audio-quality", "speech-rhythm"}:
                audio_ids.append(item.id)
                audio_metrics.append(item.metadata)
        result.evidence.append(EvidenceUnit(
            id=f"{path}:window:{start}",
            file_name=path.name,
            modality="video-window",
            locator=f"{prefix}{path.name}#t={start:.0f}-{end:.0f}",
            text=" ".join(transcript)[:3000],
            metadata={
                "startSeconds": start,
                "endSeconds": end,
                "evidenceIds": evidence_ids,
                "frameEvidenceIds": frame_ids,
                "audioEvidenceIds": audio_ids,
                "visualSummary": visual_summaries,
                "audioMetrics": audio_metrics,
            },
        ))


def _upload_artifact(path: Path) -> str | None:
    try:
        from app.config import get_settings
        from minio import Minio
        settings = get_settings()
        if not settings.minio_endpoint or not settings.minio_access_key or not settings.minio_secret_key:
            return None
        client = Minio(settings.minio_endpoint.replace("http://", "").replace("https://", ""), access_key=settings.minio_access_key, secret_key=settings.minio_secret_key, secure=settings.minio_secure)
        analysis_path = _compressed_analysis_image(path)
        key = f"analysis-artifacts/frames/{uuid4()}-{analysis_path.name}"
        client.fput_object(settings.minio_bucket, key, str(analysis_path), content_type="image/jpeg")
        return key
    except Exception:
        return None


def _compressed_analysis_image(path: Path) -> Path:
    """Create a bounded visual-analysis copy; originals remain untouched in MinIO."""
    try:
        from PIL import Image
        target = path.with_name(f"{path.stem}-analysis.jpg")
        with Image.open(path) as image:
            image = image.convert("RGB")
            image.thumbnail((1600, 1600))
            image.save(target, "JPEG", quality=78, optimize=True)
        return target
    except Exception:
        return path
